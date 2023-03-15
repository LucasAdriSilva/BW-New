from flask import Blueprint, session, request
from methods import getUrl, deletpath, countFree, send_email_with_anexo
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from flask_pymongo import PyMongo
from model.user import User
from model.db import Db
import shutil, os

mongo = PyMongo()

feature = Blueprint('feature', __name__, template_folder='templates', static_folder='static')

@feature.route("/convert")
def conversor():
  if "username" in session:
    deletpath()
    # resValidation = countFree()
    # if resValidation:
    #   return getUrl(
    #     "signature.html", prop="msgFal",
    #     value="Você atingiu o número máximo de conversões, pra continuar compre mais converções"
    #   )
    # else:
    return getUrl("conversor.html")
  else:
    return getUrl("login.html")

@feature.route('/autenticar', methods=['GET'])
def autenticar():
  if "username" in session:
    vsl = request.args.get('vsl')
    vsl = vsl.split('\n')
    Db.save_vsl(vsl)

    email = request.args.get('email')
    email = email.split(',')
    x = []

    # Faz o tratamento da string enviada pelo front
    # Add em um array e se estiver espaço vazio é excluido automaticamente
    for item in vsl:
      st = item.replace('\r', '')
      st = st.strip()
      x.append(st)

    # Remove linhas vazias
    while ("" in x):
      x.remove("")

    try:
      # Abrindo PowerPoint
      apresentacao = Presentation()
      apresentacao.slide_width = Inches(16)
      apresentacao.slide_height = Inches(9)

      # Cada item no array ira fazer esse bloco
      for line in x:

        # Criando um slide
        slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

        # Criando e centralizando a TextBox
        x = Inches(3)
        y = Inches(3.5)
        largura = Inches(10)
        altura = Inches(2)
        caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)
        caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)

        # Formatando o TextBox
        caixa_texto.text = line
        caixa_texto.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        caixa_texto.text_frame.paragraphs[0].font.size = Pt(40)
        caixa_texto.text_frame.paragraphs[0].font.name = "Arial"
        caixa_texto.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        caixa_texto.text_frame.word_wrap = True

      # Salvando o arquivo pptx
      apresentacao.save("VSL.pptx")

      # Salvando em um diretorio proprio da sessão
      name = session["username"]
      dir = f"routes/feature/static/{name}/"
      os.makedirs(dir)

      # Alterando o arquivo para um diretorio unico da sessão
      diretorio = f"routes/feature/static/{name}/VSL.pptx"
      shutil.move('VSL.pptx', diretorio)

      if email != None:
        send_email_with_anexo(
          email,
          "Envio automatico da VSL, disponviel para download em anexo desse e-mail",
          "Envio da VSL gerada", diretorio)

      user_found = Db.get_name(session["username"])
      res = user_found.data['count'] - 1
    
      Db.update_user(session["username"],{"count": res})

    except Exception as e:
      print(e)
  else:
    return getUrl("login.html")

  return getUrl("resultado.html", prop="dir", value=f"static/{name}/VSL.pptx")

@feature.route('/add')
def add():
  if "username" in session: 
    Db.update_user(session['username'], {'count': 10})
    return getUrl('conversor.html')
  else:
    return getUrl("login.html")